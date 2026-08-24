"""Create an editable PowerPoint interview presentation for this capstone."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "AI_Security_Capstone_Interview_Presentation.pptx"

NAVY = RGBColor(11, 27, 50)
BLUE = RGBColor(26, 90, 169)
CYAN = RGBColor(35, 183, 202)
MINT = RGBColor(72, 201, 176)
ORANGE = RGBColor(255, 165, 75)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(28, 42, 58)
MUTED = RGBColor(91, 108, 127)
PALE = RGBColor(242, 247, 251)
PALE_BLUE = RGBColor(228, 240, 252)
LINE = RGBColor(211, 222, 232)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def shape(slide, kind, x, y, w, h, fill, line=None, radius=True):
    item = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    item.fill.solid()
    item.fill.fore_color.rgb = fill
    item.line.color.rgb = line or fill
    return item


def add_background(slide, number, title, subtitle=None):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.22, CYAN)
    add_text(slide, title, 0.68, 0.48, 11.7, 0.5, 27, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.7, 1.03, 11.3, 0.34, 11, MUTED)
    add_text(slide, "AI SECURITY CAPSTONE", 0.7, 7.08, 2.7, 0.18, 8, MUTED, True)
    add_text(slide, f"{number:02d}", 12.08, 7.01, 0.55, 0.26, 10, BLUE, True, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, font_size=17, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = "• " + item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(11)
    return box


def add_card(slide, title, body, x, y, w, h, accent=BLUE):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, PALE, LINE)
    shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, accent)
    add_text(slide, title, x + 0.28, y + 0.24, w - 0.5, 0.34, 16, NAVY, True)
    add_text(slide, body, x + 0.28, y + 0.72, w - 0.5, h - 0.9, 12.5, MUTED)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    shape(slide, MSO_SHAPE.ARC, 8.55, -0.85, 5.9, 5.9, BLUE)
    shape(slide, MSO_SHAPE.ARC, 9.4, 0.0, 4.0, 4.0, CYAN)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, 0.7, 2.2, 0.38, CYAN, CYAN)
    add_text(slide, "INTERVIEW PRESENTATION", 0.88, 0.77, 1.85, 0.18, 9, NAVY, True)
    add_text(slide, "AI Security\nCapstone", 0.75, 1.5, 7.2, 1.55, 39, WHITE, True)
    add_text(slide, "A local, evidence-driven journey from LLM risk assessment\nto secure engineering and validation.", 0.8, 3.3, 6.8, 0.7, 18, RGBColor(207, 224, 241))
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.78, 5.0, 5.85, 0.95, RGBColor(21, 49, 81), RGBColor(21, 49, 81))
    add_text(slide, "OWASP LLM  •  MITRE ATLAS  •  MCP  •  DevSecOps", 1.05, 5.33, 5.3, 0.25, 14, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "Prepared for interview discussion", 0.8, 6.78, 3.6, 0.2, 10, RGBColor(165, 194, 220))


def add_project_slide(prs, number, title, objective, implementation, evidence, takeaway, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, number, title, "Project summary, implementation, evidence, and interview message")
    add_card(slide, "OBJECTIVE", objective, 0.7, 1.55, 3.95, 2.0, accent)
    add_card(slide, "IMPLEMENTATION", implementation, 4.82, 1.55, 3.8, 2.0, accent)
    add_card(slide, "EVIDENCE", evidence, 8.8, 1.55, 3.8, 2.0, accent)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 4.05, 11.9, 1.55, PALE_BLUE, PALE_BLUE)
    add_text(slide, "INTERVIEW TAKEAWAY", 1.0, 4.38, 2.3, 0.25, 11, BLUE, True)
    add_text(slide, takeaway, 1.0, 4.72, 11.0, 0.54, 17, NAVY, True)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 2, "The story of the capstone", "A practical security lifecycle for an LLM-enabled application")
    stages = [
        ("01", "Assess", "OWASP LLM risk baseline"),
        ("02", "Defend", "Input and output controls"),
        ("03", "Test", "Red team and ATLAS mapping"),
        ("04", "Extend", "MCP tools and resources"),
        ("05", "Automate", "CI security gates"),
        ("06", "Validate", "Scoped pen-test methodology"),
    ]
    for i, (num, name, desc) in enumerate(stages):
        x = 0.7 + (i % 3) * 4.05
        y = 1.65 + (i // 3) * 2.12
        shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 3.6, 1.58, PALE, LINE)
        add_text(slide, num, x + 0.25, y + 0.25, 0.55, 0.36, 17, CYAN, True)
        add_text(slide, name, x + 0.9, y + 0.23, 2.3, 0.32, 18, NAVY, True)
        add_text(slide, desc, x + 0.9, y + 0.72, 2.3, 0.42, 12, MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 3, "Security model: trust boundaries", "The model is one component—not the enforcement point")
    boundaries = [
        ("User input", "May contain prompt injection", ORANGE),
        ("Retrieved content", "May be poisoned or untrusted", ORANGE),
        ("Model output", "May be inaccurate or unsafe", ORANGE),
        ("Tool arguments", "May trigger real actions", ORANGE),
        ("Dependencies", "May introduce supply-chain risk", ORANGE),
    ]
    for i, (name, desc, accent) in enumerate(boundaries):
        x = 0.72 + i * 2.48
        shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.75, 2.05, 2.0, PALE, LINE)
        shape(slide, MSO_SHAPE.OVAL, x + 0.73, 1.98, 0.58, 0.58, accent, accent)
        add_text(slide, "!", x + 0.91, 2.09, 0.2, 0.2, 16, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(slide, name, x + 0.16, 2.78, 1.72, 0.28, 13, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.15, 3.15, 1.75, 0.42, 10, MUTED, align=PP_ALIGN.CENTER)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.35, 4.55, 10.6, 0.95, NAVY, NAVY)
    add_text(slide, "Validate at the boundary  •  Enforce least privilege  •  Log important decisions", 1.62, 4.89, 10.1, 0.27, 16, WHITE, True, align=PP_ALIGN.CENTER)

    add_project_slide(prs, 4, "Project 1 | OWASP LLM risk baseline",
                      "Assess an intentionally vulnerable Flask chat application against common LLM risks.",
                      "Flask on 127.0.0.1:5000 calls local Ollama. 16 KiB request cap and an in-memory 5 requests / 60 seconds rate limit.",
                      "OWASP assessment, local lab data, and source review.",
                      "A system prompt and basic rate limits are useful—but production security also requires access control, output validation, monitoring, and human approval.", BLUE)
    add_project_slide(prs, 5, "Project 2 | Prompt-injection defense",
                      "Add practical guardrails around the LLM endpoint.",
                      "POST /ask on 127.0.0.1:5001. 500-character limit, injection/Base64 regex checks, and output scanning for secrets and URLs.",
                      "Source code in app.py and validate.py; documented tested payloads.",
                      "Regex is a fast first layer, not a complete solution. Explain its bypass and false-positive tradeoffs, then describe defense in depth.", MINT)
    add_project_slide(prs, 6, "Project 3 | Red team + MITRE ATLAS",
                      "Run authorized adversarial prompts and map behavior to an AI-security framework.",
                      "20 prompts across injection, jailbreak, role-play, and encoded payloads; five ATLAS-aligned defensive tests.",
                      "Timestamped JSON results: redteam_results.json and atlas_results.json.",
                      "A model response is evidence to review—not automatic proof of compromise. Define safe success criteria and record reproducible context.", CYAN)
    add_project_slide(prs, 7, "Project 4 | MCP security assessment",
                      "Demonstrate tool injection, resource poisoning, and capability chaining with harmless local fixtures.",
                      "A restricted demo tool, a poisoned resource fixture, and server-side invocation logs. All activity stays local.",
                      "Attack JSON results plus fixture-server logs provide prompt, tool, resource, and response evidence.",
                      "Never let model text directly authorize consequential actions: confirm intent, validate server-side, scope privileges, and log every tool call.", ORANGE)
    add_project_slide(prs, 8, "Project 5 | DevSecOps security gates",
                      "Shift security left by running checks automatically on every push and pull request.",
                      "Sequential GitHub Actions gates: Bandit SAST → Safety dependencies → Gitleaks → custom prompt-injection scanner.",
                      "Workflow YAML, local validation script, pipeline design, and controlled fail/fix plan.",
                      "CI gives early feedback and blocks failed gates, but it complements—not replaces—threat modeling, review, and runtime controls.", MINT)
    add_project_slide(prs, 9, "Project 6 | Pen-test methodology",
                      "Prepare a scoped, evidence-led assessment of the local AI app and local DVWA.",
                      "Reconnaissance, source review, ZAP/Burp preparation, findings tracker, and CVSS process.",
                      "Raw Nmap/Dirb/Nikto/ZAP artifacts; findings remain unconfirmed until manually reproduced.",
                      "Do not confuse a scanner alert with a vulnerability. Confirm scope, reproduce safely, preserve evidence, then rate impact.", BLUE)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 10, "MCP: the interview deep dive", "Why tools and resources are a different class of LLM risk")
    add_card(slide, "TOOL INJECTION", "Untrusted text can influence a model to request a restricted tool. The local fixture recorded the tool call and harmless response.", 0.7, 1.55, 3.8, 2.25, ORANGE)
    add_card(slide, "RESOURCE POISONING", "A resource carried hidden instructions. The model response demonstrated a behavior change in a controlled local test.", 4.77, 1.55, 3.8, 2.25, ORANGE)
    add_card(slide, "CONTROL PATTERN", "User confirmation • allowlisted tools • narrow schemas • server validation • provenance • scoped credentials • audit logs", 8.84, 1.55, 3.8, 2.25, MINT)
    add_text(slide, "Strong sound bite: “The model can suggest an action; policy-enforcing application code must decide whether the action is allowed.”", 1.0, 4.45, 11.2, 0.6, 18, NAVY, True, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 11, "Defense in depth", "Controls should work together across the request-to-action path")
    controls = [
        ("1", "Input", "Size limits, injection checks, content classification"),
        ("2", "Context", "Separate instructions from untrusted retrieved data"),
        ("3", "Model", "Constrain role and request structured results"),
        ("4", "Output", "Scan, encode, validate, and require review"),
        ("5", "Tools", "Authorize, allowlist, validate arguments, audit"),
    ]
    for i, (num, name, desc) in enumerate(controls):
        y = 1.48 + i * 0.92
        shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.86, y, 0.56, 0.56, BLUE, BLUE)
        add_text(slide, num, 1.03, y + 0.14, 0.2, 0.2, 14, WHITE, True, align=PP_ALIGN.CENTER)
        add_text(slide, name, 1.7, y + 0.08, 1.3, 0.27, 15, NAVY, True)
        add_text(slide, desc, 3.0, y + 0.09, 7.8, 0.33, 13, MUTED)
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.85, 6.22, 11.55, 0.47, PALE_BLUE, PALE_BLUE)
    add_text(slide, "No single control is sufficient: one failed layer should not become an incident.", 1.1, 6.34, 11.0, 0.2, 13, BLUE, True, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 12, "Evidence and engineering discipline", "How the work stays credible")
    add_bullets(slide, [
        "Authorized local scope: no external targets, destructive activity, or real credentials.",
        "Repeatable records: JSON results, server-side logs, source files, workflow YAML, and reports.",
        "Clear labels: distinguish test observations, scanner alerts, and manually confirmed findings.",
        "Honest limitations: model behavior varies; simple regex defenses and local fixtures do not prove production security.",
        "Actionable output: every observed risk is linked to a control—least privilege, validation, provenance, logging, or human review.",
    ], 0.95, 1.55, 11.45, 4.8, 17)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 13, "Questions I am ready to answer", "Use these prompts to rehearse before the interview")
    questions = [
        "Why is a system prompt not an authorization boundary?",
        "How does prompt injection differ from SQL injection?",
        "What makes an MCP tool call risky, and how would you secure it?",
        "What are the limits of regex-based prompt filtering?",
        "How do you turn a red-team observation into a defensible finding?",
        "Why should a scanner alert not receive a CVSS score immediately?",
    ]
    for i, q in enumerate(questions):
        x = 0.75 + (i % 2) * 6.15
        y = 1.55 + (i // 2) * 1.42
        shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 5.7, 0.94, PALE, LINE)
        add_text(slide, q, x + 0.28, y + 0.23, 5.12, 0.44, 13, NAVY, True)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    shape(slide, MSO_SHAPE.ARC, 8.6, 2.25, 5.4, 5.4, BLUE)
    add_text(slide, "Thank you", 0.8, 1.35, 6.4, 0.65, 35, WHITE, True)
    add_text(slide, "Security is not a prompt—it is a system of\ntrust boundaries, controls, evidence, and review.", 0.84, 2.25, 7.0, 0.85, 20, RGBColor(206, 225, 242))
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.82, 4.2, 6.2, 1.0, RGBColor(21, 49, 81), RGBColor(21, 49, 81))
    add_text(slide, "Questions & discussion", 1.08, 4.54, 5.7, 0.26, 18, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "AI Security Capstone | Local, authorized, evidence-driven", 0.84, 6.75, 5.3, 0.2, 10, RGBColor(165, 194, 220))

    prs.core_properties.title = "AI Security Capstone — Interview Presentation"
    prs.core_properties.subject = "Interview presentation covering Projects 1–6"
    prs.core_properties.author = "AI Security Capstone"
    prs.save(OUTPUT)
    print(f"Created {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
